import os
import shutil
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# --- 配色方案 (VI 规范) ---
C_PRIMARY = RGBColor(0, 84, 166)    # 主色：Sixhill Blue
C_SECONDARY = RGBColor(240, 242, 245) # 辅色：浅灰背景
C_ACCENT = RGBColor(255, 107, 53)   # 强调色：活力橙
C_TEXT_DARK = RGBColor(51, 51, 51)  # 正文深灰
C_TEXT_LIGHT = RGBColor(128, 128, 128) # 注释浅灰

def add_beautified_slide(prs, title_text, layout_idx=5):
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    
    # 1. 顶部企业VI装饰条
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.15))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = C_PRIMARY
    top_bar.line.fill.background()

    # 2. 统一标题层级与对齐
    title_shape = slide.shapes.title
    title_shape.left = Inches(0.5)
    title_shape.top = Inches(0.4)
    title_shape.width = Inches(9)
    title_shape.height = Inches(0.8)
    title_shape.text = title_text
    
    tf = title_shape.text_frame
    p = tf.paragraphs[0]
    p.font.name = 'Microsoft YaHei'
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = C_PRIMARY
    p.alignment = PP_ALIGN.LEFT
    
    # 3. 底部页脚留白与辅助网格线视觉引导
    footer = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(9), Inches(0.4))
    footer.text_frame.text = "Sixhill Confidential | Data-Driven Weekly Report"
    footer.text_frame.paragraphs[0].font.size = Pt(10)
    footer.text_frame.paragraphs[0].font.color.rgb = C_TEXT_LIGHT
    footer.text_frame.paragraphs[0].font.name = 'Arial'
    
    return slide

def add_placeholder_box(slide, left, top, width, height, text, is_accent=False):
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(255, 255, 255)
    box.line.color.rgb = C_ACCENT if is_accent else C_PRIMARY
    box.text_frame.text = text
    for p in box.text_frame.paragraphs:
        p.font.color.rgb = C_TEXT_DARK
        p.font.size = Pt(14)
        p.font.name = 'Microsoft YaHei'
    return box

def create_beautified_ppt(filename):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 0: 封面
    slide0 = prs.slides.add_slide(prs.slide_layouts[6]) # Blank
    bg_shape = slide0.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = C_PRIMARY
    bg_shape.line.fill.background()
    
    title_box = slide0.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
    tf = title_box.text_frame
    p = tf.add_paragraph()
    p.text = "Sixhill 数据驱动型周汇报"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.font.name = 'Microsoft YaHei'
    
    p2 = tf.add_paragraph()
    p2.text = "Automated Business Performance Review"
    p2.font.size = Pt(20)
    p2.font.color.rgb = C_ACCENT
    p2.font.name = 'Arial'

    # Slide 1: 业绩汇总
    slide1 = add_beautified_slide(prs, "1. 上周全部分销渠道周业绩汇总")
    add_placeholder_box(slide1, 0.5, 1.5, 5.5, 4.8, "[此处嵌入 OLE 链接: Excel 表格数据区]\n支持 VBA 一键刷新")
    add_placeholder_box(slide1, 6.2, 1.5, 3.3, 2.3, "[此处嵌入 OLE 链接: Excel 动态条形图]", True)
    add_placeholder_box(slide1, 6.2, 4.0, 3.3, 2.3, "[此处嵌入 OLE 链接: Excel 环形图]", True)

    # Slide 2: 各渠道具体完成情况与数据分析
    slide2 = add_beautified_slide(prs, "2. 各渠道具体完成情况与数据分析")
    # 4 KPI cards
    for i in range(4):
        kpi = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5 + i*2.3), Inches(1.5), Inches(2.1), Inches(1))
        kpi.fill.solid()
        kpi.fill.fore_color.rgb = C_SECONDARY
        kpi.line.fill.background()
        kpi.text_frame.text = f"KPI {i+1}\n(销售额/完成率/同环比)"
        for p in kpi.text_frame.paragraphs:
            p.font.color.rgb = C_PRIMARY
            p.font.size = Pt(12)
            p.font.bold = True
            
    add_placeholder_box(slide2, 0.5, 2.8, 9.0, 2.5, "[中间留空: 链接 Excel 折线图]", True)
    add_placeholder_box(slide2, 0.5, 5.5, 4.3, 1.2, "异常说明:\n[在此填写异常原因]")
    add_placeholder_box(slide2, 5.2, 5.5, 4.3, 1.2, "结论建议:\n[在此填写改进建议]")

    # Slide 3: 本周工作完成情况说明
    slide3 = add_beautified_slide(prs, "3. 各渠道本周工作完成情况说明")
    add_placeholder_box(slide3, 0.5, 1.5, 9.0, 3.8, "[此处嵌入 Excel 4栏表格: 渠道/工作项/状态/备注]")
    add_placeholder_box(slide3, 0.5, 5.5, 9.0, 1.2, "共性风险总结:\n[在此填写各渠道共性风险与应对措施]")

    # Slide 4: 下周工作安排说明
    slide4 = add_beautified_slide(prs, "4. 各渠道下周工作安排说明")
    date_box = slide4.shapes.add_textbox(Inches(6.5), Inches(0.4), Inches(3), Inches(0.5))
    date_box.text_frame.text = "周期: [链接自 Excel 日期]"
    date_box.text_frame.paragraphs[0].font.color.rgb = C_ACCENT
    date_box.text_frame.paragraphs[0].font.size = Pt(14)
    date_box.text_frame.paragraphs[0].font.bold = True
    
    add_placeholder_box(slide4, 0.5, 1.5, 4.3, 5.0, "左侧重点事项编号:\n1.\n2.\n3.")
    add_placeholder_box(slide4, 5.2, 1.5, 4.3, 5.0, "右侧预期输出与优先级(高/中/低):\n- 输出A (高)\n- 输出B (中)")

    # Slide 5: 剩余目标额说明
    slide5 = add_beautified_slide(prs, "5. 未来一周及本月剩余目标额说明")
    add_placeholder_box(slide5, 0.5, 1.5, 9.0, 3.8, "[此处嵌入 Excel 表格及数据进度条(Data Bar)]")
    add_placeholder_box(slide5, 0.5, 5.5, 9.0, 1.2, "资源需求与风险:\n[在此填写需协调的资源与潜在风险]")

    # Slide 6: AI 驱动的业绩增长方案
    slide6 = add_beautified_slide(prs, "6. AI 驱动的业绩增长方案与计划")
    add_placeholder_box(slide6, 0.5, 1.5, 9.0, 1.0, "机会点洞察:\n[填入 AI 分析得出的增长结论]")
    add_placeholder_box(slide6, 0.5, 2.7, 5.5, 2.8, "[行动方案 5 列表格:\n动作 | 负责人 | 起止日期 | 预期提升]")
    add_placeholder_box(slide6, 6.2, 2.7, 3.3, 2.8, "[右侧预留甘特图示意区]", True)
    add_placeholder_box(slide6, 0.5, 5.7, 9.0, 1.0, "KPI 成功指标区:\n[如: 转化率提升 5%, 客单价提高 10%]")

    prs.save(filename)

if __name__ == "__main__":
    create_beautified_ppt("Sixhill_Beautified_Report.pptx")
    shutil.copy("Sixhill_Beautified_Report.pptx", "Sixhill_Beautified_Master.potx")
    print("Full 6-slide V2 Assets generated successfully.")
