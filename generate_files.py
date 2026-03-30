import os
import xlsxwriter
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from fpdf import FPDF

SIXHILL_BLUE = '#0054A6'
SIXHILL_RGB = RGBColor(0, 84, 166)

def create_excel(filename):
    workbook = xlsxwriter.Workbook(filename)
    
    # Formats
    header_format = workbook.add_format({
        'bold': True, 'bg_color': SIXHILL_BLUE, 'font_color': 'white', 
        'border': 1, 'align': 'center', 'valign': 'vcenter'
    })
    cell_format = workbook.add_format({'border': 1, 'align': 'center', 'valign': 'vcenter'})
    pct_format = workbook.add_format({'border': 1, 'align': 'center', 'valign': 'vcenter', 'num_format': '0.00%'})
    money_format = workbook.add_format({'border': 1, 'align': 'center', 'valign': 'vcenter', 'num_format': '#,##0'})
    
    # Sheet 1: 业绩汇总
    ws1 = workbook.add_worksheet('上周业绩汇总')
    headers1 = ['渠道', '本周目标', '本周实际', '本周完成率', '本周环比', '本月目标', '本月实际', '本月完成率']
    for col_num, data in enumerate(headers1):
        ws1.write(0, col_num, data, header_format)
    
    channels = ['LSH', 'MBJ', 'Natalia', 'Phoenix', '其他']
    row = 1
    for ch in channels:
        ws1.write(row, 0, ch, cell_format)
        ws1.write(row, 1, 100000, money_format) # 本周目标
        ws1.write(row, 2, 85000, money_format)  # 本周实际
        ws1.write_formula(row, 3, f'=C{row+1}/B{row+1}', pct_format) # 本周完成率
        ws1.write(row, 4, 0.05, pct_format)     # 本周环比
        ws1.write(row, 5, 400000, money_format) # 本月目标
        ws1.write(row, 6, 320000, money_format) # 本月实际
        ws1.write_formula(row, 7, f'=G{row+1}/F{row+1}', pct_format) # 本月完成率
        row += 1
        
    # 整体汇总
    ws1.write(row, 0, '整体汇总', header_format)
    ws1.write_formula(row, 1, f'=SUM(B2:B{row})', money_format)
    ws1.write_formula(row, 2, f'=SUM(C2:C{row})', money_format)
    ws1.write_formula(row, 3, f'=C{row+1}/B{row+1}', pct_format)
    ws1.write(row, 4, 0.08, pct_format)
    ws1.write_formula(row, 5, f'=SUM(F2:F{row})', money_format)
    ws1.write_formula(row, 6, f'=SUM(G2:G{row})', money_format)
    ws1.write_formula(row, 7, f'=G{row+1}/F{row+1}', pct_format)
    
    # Conditional formatting for completion rate (<80% red, >=100% green)
    red_format = workbook.add_format({'font_color': '#9C0006', 'bg_color': '#FFC7CE'})
    green_format = workbook.add_format({'font_color': '#006100', 'bg_color': '#C6EFCE'})
    
    ws1.conditional_format(f'D2:D{row+1}', {'type': 'cell', 'criteria': '<', 'value': 0.8, 'format': red_format})
    ws1.conditional_format(f'D2:D{row+1}', {'type': 'cell', 'criteria': '>=', 'value': 1.0, 'format': green_format})
    ws1.conditional_format(f'H2:H{row+1}', {'type': 'cell', 'criteria': '<', 'value': 0.8, 'format': red_format})
    ws1.conditional_format(f'H2:H{row+1}', {'type': 'cell', 'criteria': '>=', 'value': 1.0, 'format': green_format})
    
    ws1.set_column('A:H', 12)

    # Sheet 2: 工作完成情况
    ws2 = workbook.add_worksheet('工作完成情况')
    headers2 = ['渠道', '工作项', '状态', '备注']
    for col_num, data in enumerate(headers2):
        ws2.write(0, col_num, data, header_format)
        
    for i in range(1, 10):
        ws2.write(i, 0, 'LSH', cell_format)
        ws2.write(i, 1, f'工作项 {i}', cell_format)
        ws2.write(i, 2, '进行中', cell_format)
        ws2.write(i, 3, '', cell_format)
        
    ws2.data_validation('C2:C20', {'validate': 'list', 'source': ['已完成', '进行中', '未开始']})
    ws2.set_column('A:A', 15)
    ws2.set_column('B:B', 30)
    ws2.set_column('C:C', 15)
    ws2.set_column('D:D', 30)

    # Sheet 3: 目标与概率
    ws3 = workbook.add_worksheet('目标与概率')
    headers3 = ['渠道', '剩余目标', '本月剩余天数', '日均需完成', '预计达成概率']
    for col_num, data in enumerate(headers3):
        ws3.write(0, col_num, data, header_format)
        
    row = 1
    for ch in channels:
        ws3.write(row, 0, ch, cell_format)
        ws3.write(row, 1, 15000, money_format)
        ws3.write(row, 2, 10, cell_format)
        ws3.write_formula(row, 3, f'=B{row+1}/C{row+1}', money_format)
        ws3.write(row, 4, 0.85, pct_format)
        row += 1
        
    ws3.conditional_format(f'E2:E{row}', {'type': 'data_bar', 'bar_color': SIXHILL_BLUE})
    ws3.set_column('A:E', 15)
    
    # Sheet 4: PPT Data (Dates, etc.)
    ws4 = workbook.add_worksheet('PPT_Data')
    ws4.write(0, 0, '本周开始日期', header_format)
    ws4.write(0, 1, '本周结束日期', header_format)
    ws4.write(1, 0, '2023-10-16', cell_format)
    ws4.write(1, 1, '2023-10-22', cell_format)
    
    workbook.close()

def add_title(slide, title_text):
    title = slide.shapes.title
    title.text = title_text
    title.text_frame.paragraphs[0].font.color.rgb = SIXHILL_RGB
    title.text_frame.paragraphs[0].font.bold = True

def create_ppt(filename):
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[5] # Title only
    
    # Slide 1
    slide1 = prs.slides.add_slide(blank_slide_layout)
    add_title(slide1, "1. 上周全部分销渠道周业绩汇总")
    # Placeholders
    txBox = slide1.shapes.add_textbox(Inches(0.5), Inches(2), Inches(4), Inches(4))
    tf = txBox.text_frame
    tf.text = "[此处使用“选择性粘贴 -> 链接”嵌入 Excel 表格区域]"
    
    txBox2 = slide1.shapes.add_textbox(Inches(5), Inches(2), Inches(4), Inches(2))
    txBox2.text_frame.text = "[右侧区域: 链接 Excel 动态条形图]"
    txBox3 = slide1.shapes.add_textbox(Inches(5), Inches(4.5), Inches(4), Inches(2))
    txBox3.text_frame.text = "[右侧区域: 链接 Excel 环形图]"

    # Slide 2
    slide2 = prs.slides.add_slide(blank_slide_layout)
    add_title(slide2, "2. 各渠道具体完成情况与数据分析")
    # 4 KPI cards
    for i in range(4):
        shape = slide2.shapes.add_shape(1, Inches(0.5 + i*2.3), Inches(1.5), Inches(2), Inches(1))
        shape.fill.solid()
        shape.fill.fore_color.rgb = SIXHILL_RGB
        shape.text = f"KPI {i+1}\n(销售额/完成率/同环比)"
        
    txBox_chart = slide2.shapes.add_textbox(Inches(0.5), Inches(3), Inches(8.5), Inches(2.5))
    txBox_chart.text_frame.text = "[中间留空: 链接 Excel 折线图]"
    
    txBox_anomaly = slide2.shapes.add_textbox(Inches(0.5), Inches(6), Inches(4), Inches(1))
    txBox_anomaly.text_frame.text = "异常说明:\n[在此填写异常原因]"
    
    txBox_advice = slide2.shapes.add_textbox(Inches(5), Inches(6), Inches(4), Inches(1))
    txBox_advice.text_frame.text = "结论建议:\n[在此填写改进建议]"

    # Slide 3
    slide3 = prs.slides.add_slide(blank_slide_layout)
    add_title(slide3, "3. 各渠道本周工作完成情况说明")
    txBox_table = slide3.shapes.add_textbox(Inches(0.5), Inches(2), Inches(8.5), Inches(3.5))
    txBox_table.text_frame.text = "[此处嵌入 Excel 4栏表格: 渠道/工作项/状态/备注]"
    
    txBox_risk = slide3.shapes.add_textbox(Inches(0.5), Inches(6), Inches(8.5), Inches(1))
    txBox_risk.text_frame.text = "共性风险总结:\n[在此填写各渠道共性风险与应对措施]"

    # Slide 4
    slide4 = prs.slides.add_slide(blank_slide_layout)
    add_title(slide4, "4. 各渠道下周工作安排说明")
    txBox_date = slide4.shapes.add_textbox(Inches(6), Inches(0.5), Inches(3), Inches(0.5))
    txBox_date.text_frame.text = "[日期区域: 链接自 Excel (开始 至 结束)]"
    
    txBox_left = slide4.shapes.add_textbox(Inches(0.5), Inches(2), Inches(4), Inches(4))
    txBox_left.text_frame.text = "左侧重点事项编号:\n1.\n2.\n3."
    
    txBox_right = slide4.shapes.add_textbox(Inches(5), Inches(2), Inches(4), Inches(4))
    txBox_right.text_frame.text = "右侧预期输出与优先级(高/中/低):\n- 输出A (高)\n- 输出B (中)"

    # Slide 5
    slide5 = prs.slides.add_slide(blank_slide_layout)
    add_title(slide5, "5. 未来一周及本月剩余目标额说明")
    txBox_target = slide5.shapes.add_textbox(Inches(0.5), Inches(2), Inches(8.5), Inches(3.5))
    txBox_target.text_frame.text = "[此处嵌入 Excel 表格及数据进度条]"
    
    txBox_req = slide5.shapes.add_textbox(Inches(0.5), Inches(6), Inches(8.5), Inches(1))
    txBox_req.text_frame.text = "资源需求与风险:\n[在此填写需协调的资源与潜在风险]"

    # Slide 6
    slide6 = prs.slides.add_slide(blank_slide_layout)
    add_title(slide6, "6. AI 驱动的业绩增长方案与计划")
    txBox_ai = slide6.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(8.5), Inches(1))
    txBox_ai.text_frame.text = "机会点洞察:\n[填入 AI 分析得出的增长结论]"
    
    txBox_action = slide6.shapes.add_textbox(Inches(0.5), Inches(3), Inches(4.5), Inches(2.5))
    txBox_action.text_frame.text = "[行动方案 5 列表格:\n动作 | 负责人 | 起止日期 | 预期提升]"
    
    txBox_gantt = slide6.shapes.add_textbox(Inches(5.5), Inches(3), Inches(3.5), Inches(2.5))
    txBox_gantt.text_frame.text = "[右侧预留甘特图示意区]"
    
    txBox_kpi = slide6.shapes.add_textbox(Inches(0.5), Inches(6), Inches(8.5), Inches(1))
    txBox_kpi.text_frame.text = "KPI 成功指标区:\n[如: 转化率提升 5%, 客单价提高 10%]"

    prs.save(filename)

def create_pdf(filename):
    pdf = FPDF()
    pdf.add_page()
    
    # Try to load a Chinese font. If not found, use Arial but English text.
    font_path = "C:\\Windows\\Fonts\\simhei.ttf"
    if os.path.exists(font_path):
        pdf.add_font("SimHei", "", font_path, uni=True)
        pdf.set_font("SimHei", size=16)
        
        pdf.set_text_color(0, 84, 166) # Sixhill Blue
        pdf.cell(200, 15, txt="Sixhill 周汇报自动化更新指南 (User Guide)", ln=True, align='C')
        
        pdf.set_font("SimHei", size=12)
        pdf.set_text_color(0, 0, 0)
        
        instructions = [
            "",
            "1. 架构说明:",
            "   本模板套装采用“数据引擎(Excel) + 视觉呈现(PPT)”分离架构。",
            "   所有数据计算与图表生成均在 Excel 中完成，PPT 仅作为展示前端。",
            "",
            "2. 日常更新步骤:",
            "   步骤一: 打开 Sixhill_Weekly_Master.xlsx，更新各个 Sheet 的数据。",
            "   步骤二: 检查条件格式(标红/标绿)与计算结果是否正确。",
            "   步骤三: 保存并保留 Excel 开启状态。",
            "   步骤四: 打开 Sixhill_Weekly_Report.pptx。",
            "   步骤五: 若系统提示“是否更新链接”，请点击“更新”。PPT 中的数据会自动刷新。",
            "",
            "3. 如何建立新的数据链接 (选择性粘贴):",
            "   如果需要在 PPT 中新增 Excel 链接:",
            "   a. 在 Excel 中选中需要展示的表格区域或图表，按 Ctrl+C 复制。",
            "   b. 切换到 PPT，在对应页面点击右键，选择“选择性粘贴 (Paste Special)”。",
            "   c. 选择“粘贴链接 (Paste Link)”，并选择“Microsoft Excel 工作表 对象”。",
            "   d. 调整嵌入对象的大小与位置。",
            "",
            "4. 注意事项:",
            "   - 请确保 Excel 和 PPT 存放在相同文件夹，或者不要随意移动文件位置，以免链接失效。",
            "   - PPT 中的数据仅能通过修改 Excel 数据源来更新，切勿在 PPT 中直接双击修改数字。"
        ]
        for line in instructions:
            pdf.cell(200, 8, txt=line, ln=True)
    else:
        pdf.set_font("Arial", size=16)
        pdf.set_text_color(0, 84, 166)
        pdf.cell(200, 15, txt="Sixhill Weekly Report User Guide", ln=True, align='C')
        pdf.set_font("Arial", size=12)
        pdf.set_text_color(0, 0, 0)
        pdf.cell(200, 10, txt="Font simhei.ttf not found. English fallback.", ln=True)
        
    pdf.output(filename)

if __name__ == "__main__":
    create_excel("Sixhill_Weekly_Master.xlsx")
    create_ppt("Sixhill_Weekly_Report.pptx")
    create_pdf("User_Guide.pdf")
    print("Files generated successfully.")
